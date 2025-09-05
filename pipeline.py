# pipeline.py
# ===============================================================
# Unified data pipeline:
# - Confluence roots/pages (crawl + fetch) + optional Coveo search labels
# - Local files of mixed types (pdf/docx/txt/md/json/csv/xlsx/pptx)
# - Chunking (RecursiveCharacterTextSplitter; tuned per type)
# - Azure OpenAI embeddings via Azure AD (MSAL)
# - Pinecone upsert (re-use index if exists; else create)
# - Detailed report (words, pages, chunks per source; totals)
#
# Run:
#   python pipeline.py
#
# Place a non-secret config.json next to this file (example at bottom).
# Secrets are pulled via ARNs embedded below (no CLI / env needed).
# ===============================================================

import os
import re
import io
import json
import time
import math
import base64
import hashlib
import logging
import mimetypes
from dataclasses import dataclass
from typing import Dict, List, Tuple, Optional, Iterable

import boto3
import requests

# ---------- optional deps (all handled defensively) ----------
try:
    import fitz  # PyMuPDF
except Exception:
    fitz = None

try:
    import pdfplumber
except Exception:
    pdfplumber = None

try:
    from pdfminer.high_level import extract_text as pdfminer_extract_text
except Exception:
    pdfminer_extract_text = None

try:
    import docx  # python-docx
except Exception:
    docx = None

try:
    import pandas as pd
except Exception:
    pd = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

try:
    from langchain.text_splitter import RecursiveCharacterTextSplitter
except Exception:
    RecursiveCharacterTextSplitter = None

try:
    from langchain_community.document_loaders import ConfluenceLoader
except Exception:
    ConfluenceLoader = None

from msal import ConfidentialClientApplication
import openai

from pinecone import Pinecone, ServerlessSpec

# Some Pinecone SDK versions have different exception import paths
try:
    from pinecone.exceptions import PineconeApiException
except Exception:
    try:
        from pinecone.exceptions.exceptions import PineconeApiException
    except Exception:
        class PineconeApiException(Exception):
            pass

# ===============================================================
# Configuration for secrets (in-file) and SDK bootstrap
# ===============================================================

# --- Replace these ARNs with your real ones (these live in code; NOT in config.json) ---
PINECONE_SECRET_ARN = "arn:aws:secretsmanager:us-east-1:123456789012:secret:pinecone-REPLACE_ME"
AZURE_OPENAI_SECRET_ARN = "arn:aws:secretsmanager:us-east-1:123456789012:secret:azure-openai-REPLACE_ME"

# --- Azure OpenAI constants (live in code) ---
AZURE_TENANT_ID   = "YOUR_TENANT_ID_OR_NAME"
AZURE_OAI_BASE    = "https://YOUR-RESOURCE.openai.azure.com/"
AZURE_OAI_VERSION = "2023-05-15"
EMBEDDING_MODEL   = "text-embedding-ada-002"  # deployment name of your Azure embedding model

# Pinecone defaults
PINECONE_DIM      = 1536
PINECONE_METRIC   = "cosine"
PINECONE_REGION   = "us-east-1"

# Ensure region defaults for boto3
os.environ.setdefault("AWS_REGION", "us-east-1")
os.environ.setdefault("AWS_DEFAULT_REGION", "us-east-1")

logging.basicConfig(level=logging.INFO, format="%(asctime)s [%(levelname)s] %(message)s")


def pinecone_config() -> Pinecone:
    """Fetch Pinecone API key from Secrets Manager and return a Pinecone client."""
    sm = boto3.client("secretsmanager", region_name="us-east-1")
    resp = sm.get_secret_value(SecretId=PINECONE_SECRET_ARN)
    data = json.loads(resp["SecretString"])
    api_key = (data.get("apiKey") or "").replace('"', "").strip()
    if not api_key:
        raise RuntimeError("Pinecone apiKey missing in secret.")
    return Pinecone(api_key=api_key)


def openai_api_config() -> str:
    """Fetch Azure SP creds from Secrets Manager, get AAD token, and configure OpenAI SDK."""
    sm = boto3.client("secretsmanager", region_name="us-east-1")
    resp = sm.get_secret_value(SecretId=AZURE_OPENAI_SECRET_ARN)
    data = json.loads(resp["SecretString"])

    client_id = data.get("AzureServicePrincipalId", "").strip()
    client_secret = (data.get("Password") or "").replace('"', "").strip()
    if not client_id or not client_secret:
        raise RuntimeError("Azure OpenAI SP credentials missing in secret.")

    app = ConfidentialClientApplication(
        client_id=client_id,
        client_credential=client_secret,
        authority=f"https://login.microsoftonline.com/{AZURE_TENANT_ID}",
    )
    scopes = ["https://cognitiveservices.azure.com/.default"]
    token_result = app.acquire_token_for_client(scopes=scopes)
    if "access_token" not in token_result:
        raise RuntimeError(f"Unable to obtain Azure AD token: {token_result}")

    openai.api_type = "azure_ad"
    openai.api_key = token_result["access_token"]
    openai.api_base = AZURE_OAI_BASE
    openai.api_version = AZURE_OAI_VERSION

    return EMBEDDING_MODEL


# ===============================================================
# Utilities
# ===============================================================

def sha1(s: str) -> str:
    return hashlib.sha1(s.encode("utf-8", "ignore")).hexdigest()

def looks_like_url(s: str) -> bool:
    if not isinstance(s, str):
        return False
    return s.startswith("http://") or s.startswith("https://")

def safe_read_text(path: str) -> str:
    with open(path, "rb") as f:
        raw = f.read()
    for enc in ("utf-8", "utf-16", "latin-1"):
        try:
            return raw.decode(enc)
        except Exception:
            continue
    # fallback
    return raw.decode("utf-8", "ignore")

def count_words(text: str) -> int:
    return len([w for w in re.findall(r"\w+", text)])

def now_iso() -> str:
    return time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime())

# =============== Confluence helpers (regex + REST crawl) ===============

CONFLUENCE_PAGE_ID_RE = re.compile(r"pages/(\d+)")

def get_base_url(url: str) -> str:
    """
    A bit stricter than the user's version: keep scheme + host until '/wiki' (Atlassian Cloud),
    else just scheme+host.
    """
    m = re.match(r"^(https?://[^/]+)(/wiki)?/.*$", url)
    if m:
        return m.group(1) + (m.group(2) or "")
    raise ValueError(f"Unable to derive Confluence base from: {url}")

def get_page_id(url: str) -> str:
    m = CONFLUENCE_PAGE_ID_RE.search(url)
    if not m:
        raise ValueError(f"No page ID found in URL: {url}")
    return m.group(1)

def confluence_get_children(base_url: str, page_id: str, auth: requests.auth.HTTPBasicAuth, limit: int = 50, start: int = 0) -> Dict:
    """
    Calls the Confluence REST API for child pages of a page_id.
    """
    url = f"{base_url}/rest/api/content/{page_id}/child/page"
    params = {"limit": limit, "start": start, "expand": "ancestors"}
    r = requests.get(url, params=params, auth=auth, timeout=30)
    r.raise_for_status()
    return r.json()

def confluence_collect_descendants(root_url: str, username: str, api_token: str, max_pages: int = 150, max_depth: int = 3) -> List[str]:
    """
    BFS crawl child pages starting from root_url up to max_pages & max_depth.
    """
    base = get_base_url(root_url)
    root_id = get_page_id(root_url)
    auth = requests.auth.HTTPBasicAuth(username, api_token)

    queue: List[Tuple[str, int]] = [(root_id, 0)]
    seen_ids: set[str] = set([root_id])
    urls: List[str] = [root_url]

    while queue and len(urls) < max_pages:
        cur_id, depth = queue.pop(0)
        if depth >= max_depth:
            continue
        # paginate children
        start = 0
        while True:
            data = confluence_get_children(base, cur_id, auth, limit=50, start=start)
            results = data.get("results", [])
            if not results:
                break
            for item in results:
                cid = item.get("id")
                if cid and cid not in seen_ids:
                    seen_ids.add(cid)
                    child_url = f"{base}/pages/{cid}"
                    urls.append(child_url)
                    if len(urls) >= max_pages:
                        break
                    queue.append((cid, depth + 1))
            if len(urls) >= max_pages:
                break
            if data.get("_links", {}).get("next"):
                start += 50
            else:
                break
    return urls

# =================== Coveo label → Confluence URLs ===================

class CoveoSearch:
    def __init__(self, organization_id: str, platform_token: str):
        self.organization_id = organization_id
        self.platform_token = platform_token.strip()
        self.base_url = "https://platform.cloud.coveo.com/rest/search/v2"
        self.search_url = f"https://{organization_id}.org.coveo.com/rest/search/v2"

    def get_token(self, user_email: str) -> str:
        url = f"{self.base_url}/token"
        payload = {
            "organizationId": self.organization_id,
            "validFor": 1800000,
            "userIds": [{"name": user_email, "provider": "Email Security Provider"}],
        }
        headers = {"authorization": f"Bearer {self.platform_token}", "content-type": "application/json"}
        r = requests.post(url, json=payload, headers=headers, timeout=30)
        r.raise_for_status()
        return r.json().get("token", "")

    def search_links(self, label: str, user_token: str) -> List[str]:
        querystring = {"organizationId": self.organization_id}
        payload = {"q": f"@conflabels={label}"}
        headers = {"authorization": f"Bearer {user_token}", "content-type": "application/json"}
        r = requests.post(self.search_url, json=payload, headers=headers, params=querystring, timeout=30)
        r.raise_for_status()
        data = r.json()
        out = []
        for res in data.get("results", []):
            uri = res.get("clickUri")
            if uri and "atlassian.net/wiki" in uri:
                out.append(uri)
        return out

# ===============================================================
# Extraction
# ===============================================================

@dataclass
class Segment:
    source_id: str          # stable id (e.g., file path or URL)
    source_type: str        # "file" | "confluence"
    locator: str            # path or URL (+ page number if applicable)
    text: str
    meta: Dict[str, str]    # any extra metadata: {"page": "3", "filetype": "pdf", ...}

def extract_confluence_pages(urls: List[str], username: str, api_token: str) -> List[Segment]:
    segments: List[Segment] = []
    if not urls:
        return segments
    if ConfluenceLoader is None:
        logging.warning("ConfluenceLoader not installed; skipping Confluence extraction.")
        return segments

    # Group by base to batch load if you want (here: page-by-page to keep metadata precise)
    for u in urls:
        try:
            base = get_base_url(u)
            pid = get_page_id(u)
            loader = ConfluenceLoader(
                url=base,
                username=username,
                api_key=api_token,
                page_ids=[pid],
                include_attachments=True
            )
            docs = loader.load()
            for d in docs:
                title = d.metadata.get("title", "")
                text = f"{title}\n\n{d.page_content}".strip()
                if not text:
                    continue
                seg = Segment(
                    source_id=u,
                    source_type="confluence",
                    locator=u,
                    text=text,
                    meta={
                        "title": title,
                        "source_url": u,
                        "type": "confluence"
                    }
                )
                segments.append(seg)
        except Exception as e:
            logging.warning(f"Confluence load failed for {u}: {e}")
    return segments

# ------------------------- File extraction -------------------------

def extract_pdf(path: str) -> List[Tuple[int, str]]:
    """
    Returns list of (page_number, text). Tries PyMuPDF -> pdfplumber -> pdfminer.
    """
    pages: List[Tuple[int, str]] = []
    # 1) PyMuPDF
    if fitz is not None:
        try:
            doc = fitz.open(path)
            for i, p in enumerate(doc, start=1):
                txt = p.get_text("text") or ""
                if txt.strip():
                    pages.append((i, txt))
            doc.close()
            if pages:
                return pages
        except Exception as e:
            logging.info(f"PyMuPDF failed on {path}: {e}")

    # 2) pdfplumber
    if pdfplumber is not None:
        try:
            with pdfplumber.open(path) as pdf:
                for i, p in enumerate(pdf.pages, start=1):
                    txt = p.extract_text() or ""
                    if txt.strip():
                        pages.append((i, txt))
            if pages:
                return pages
        except Exception as e:
            logging.info(f"pdfplumber failed on {path}: {e}")

    # 3) pdfminer (whole doc)
    if pdfminer_extract_text is not None:
        try:
            txt = pdfminer_extract_text(path) or ""
            if txt.strip():
                # No page granularity from here; return as page 1
                return [(1, txt)]
        except Exception as e:
            logging.info(f"pdfminer failed on {path}: {e}")

    return pages

def extract_docx(path: str) -> str:
    if docx is None:
        logging.warning("python-docx not installed; cannot read DOCX")
        return ""
    try:
        d = docx.Document(path)
        paras = [p.text for p in d.paragraphs if p.text and p.text.strip()]
        return "\n".join(paras)
    except Exception as e:
        logging.warning(f"DOCX read failed {path}: {e}")
        return ""

def extract_pptx(path: str) -> str:
    if Presentation is None:
        logging.warning("python-pptx not installed; cannot read PPTX")
        return ""
    try:
        prs = Presentation(path)
        out = []
        for si, slide in enumerate(prs.slides, start=1):
            buf = [f"[Slide {si}]"]
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    t = (shape.text or "").strip()
                    if t:
                        buf.append(t)
            if len(buf) > 1:
                out.append("\n".join(buf))
        return "\n\n".join(out)
    except Exception as e:
        logging.warning(f"PPTX read failed {path}: {e}")
        return ""

def extract_xlsx(path: str) -> str:
    if pd is None:
        logging.warning("pandas not installed; cannot read XLSX")
        return ""
    try:
        xl = pd.ExcelFile(path)
        out = []
        for sheet in xl.sheet_names:
            df = xl.parse(sheet)
            # limit very large sheets to avoid ballooning
            if len(df) > 1000:
                df = df.head(1000)
            csv_txt = df.to_csv(index=False)
            out.append(f"[Sheet: {sheet}]\n{csv_txt}")
        return "\n\n".join(out)
    except Exception as e:
        logging.warning(f"XLSX read failed {path}: {e}")
        return ""

def extract_csv(path: str) -> str:
    if pd is None:
        try:
            return safe_read_text(path)
        except:
            return ""
    try:
        df = pd.read_csv(path)
        if len(df) > 5000:
            df = df.head(5000)
        return df.to_csv(index=False)
    except Exception as e:
        logging.warning(f"CSV read failed {path}: {e}")
        return safe_read_text(path)

def extract_json(path: str) -> str:
    try:
        with open(path, "r", encoding="utf-8") as f:
            data = json.load(f)
        return json.dumps(data, indent=2, ensure_ascii=False)
    except Exception:
        return safe_read_text(path)

def extract_text_like(path: str) -> str:
    return safe_read_text(path)

def extract_file_segments(path: str) -> List[Segment]:
    ext = (os.path.splitext(path)[1] or "").lower()
    segments: List[Segment] = []
    if ext == ".pdf":
        pages = extract_pdf(path)
        if pages:
            for pg, txt in pages:
                segments.append(Segment(
                    source_id=path, source_type="file",
                    locator=f"{path}#page={pg}", text=txt,
                    meta={"type": "pdf", "page": str(pg), "file_path": path}
                ))
        else:
            logging.warning(f"No text from PDF: {path}")
    elif ext in (".docx",):
        txt = extract_docx(path)
        if txt.strip():
            segments.append(Segment(path, "file", path, txt, {"type": "docx", "file_path": path}))
    elif ext in (".pptx",):
        txt = extract_pptx(path)
        if txt.strip():
            segments.append(Segment(path, "file", path, txt, {"type": "pptx", "file_path": path}))
    elif ext in (".xlsx", ".xls"):
        txt = extract_xlsx(path)
        if txt.strip():
            segments.append(Segment(path, "file", path, txt, {"type": "xlsx", "file_path": path}))
    elif ext in (".csv",):
        txt = extract_csv(path)
        if txt.strip():
            segments.append(Segment(path, "file", path, txt, {"type": "csv", "file_path": path}))
    elif ext in (".json",):
        txt = extract_json(path)
        if txt.strip():
            segments.append(Segment(path, "file", path, txt, {"type": "json", "file_path": path}))
    else:
        # txt / md / xml / yaml / unknown
        txt = extract_text_like(path)
        if txt.strip():
            segments.append(Segment(path, "file", path, txt, {"type": ext.lstrip(".") or "text", "file_path": path}))
    return segments

# ===============================================================
# Chunking
# ===============================================================

def pick_chunk_params(seg_type: str) -> Tuple[int, int]:
    """
    Tune chunk sizes a bit per type.
    """
    t = (seg_type or "").lower()
    if t == "pdf":
        return (1200, 200)
    if t in ("docx", "pptx"):
        return (1000, 200)
    if t in ("xlsx", "csv", "json"):
        return (1500, 200)
    # default
    return (1000, 200)

def chunk_segments(segments: List[Segment]) -> List[Dict]:
    if not segments:
        return []
    if RecursiveCharacterTextSplitter is None:
        # very basic fallback
        out = []
        for s in segments:
            text = s.text
            for i in range(0, len(text), 1000):
                chunk = text[i:i+1000]
                out.append({
                    "id": f"{sha1(s.locator)}_{i//1000}",
                    "text": chunk,
                    "metadata": {**s.meta, "locator": s.locator, "source_type": s.source_type}
                })
        return out

    out: List[Dict] = []
    # group by params to avoid creating many splitters
    groups: Dict[Tuple[int, int], List[Segment]] = {}
    for s in segments:
        cs, ov = pick_chunk_params(s.meta.get("type", ""))
        groups.setdefault((cs, ov), []).append(s)

    for (cs, ov), segs in groups.items():
        splitter = RecursiveCharacterTextSplitter(chunk_size=cs, chunk_overlap=ov)
        for s in segs:
            chunks = splitter.split_text(s.text)
            for i, c in enumerate(chunks):
                out.append({
                    "id": f"{sha1(s.locator)}_{i}",
                    "text": c,
                    "metadata": {**s.meta, "locator": s.locator, "source_type": s.source_type}
                })
    return out

# ===============================================================
# Embeddings + Pinecone
# ===============================================================

def _list_index_names(pc: Pinecone) -> List[str]:
    try:
        existing = pc.list_indexes()
        if not existing:
            return []
        # Normalize
        if isinstance(existing[0], dict):
            return [i["name"] for i in existing]
        return list(existing)
    except Exception:
        return []

def _ensure_index(pc: Pinecone, index_name: str, dimension: int = PINECONE_DIM):
    """Describe or create index; tolerate already-exists races."""
    try:
        pc.describe_index(index_name)
        return
    except PineconeApiException as e:
        if "NOT_FOUND" in str(e) or "404" in str(e):
            pass
        else:
            # try create anyway
            pass
    except Exception:
        pass

    try:
        pc.create_index(
            name=index_name,
            dimension=dimension,
            metric=PINECONE_METRIC,
            spec=ServerlessSpec(cloud="aws", region=PINECONE_REGION),
            deletion_protection="disabled"
        )
    except PineconeApiException as e:
        if "ALREADY_EXISTS" in str(e) or "409" in str(e):
            pass
        else:
            raise
    except Exception:
        # ignore; final describe below
        pass

    pc.describe_index(index_name)

def _private_host_for_index(pc: Pinecone, index_name: str) -> str:
    desc = pc.describe_index(index_name)
    host = desc["host"]
    parts = host.split(".")
    # <name>-<project>.svc.<region>.pinecone.io -> <name>-<project>.private.<region>.pinecone.io
    return f"{'.'.join(parts[:2])}.private.{'.'.join(parts[2:])}"

def embed_texts(batch: List[str], engine: str, retries: int = 3) -> List[List[float]]:
    for attempt in range(retries):
        try:
            # Azure AD configured in openai_api_config()
            resp = openai.Embedding.create(input=batch, engine=engine)
            return [d["embedding"] for d in resp["data"]]
        except Exception as e:
            wait = 2 ** attempt
            logging.warning(f"Embedding error (attempt {attempt+1}/{retries}): {e}. Retrying in {wait}s")
            time.sleep(wait)
    return []

def upsert_chunks(pc: Pinecone, index_name: str, chunks: List[Dict], engine: str, batch_size: int = 32):
    _ensure_index(pc, index_name, dimension=PINECONE_DIM)
    host = _private_host_for_index(pc, index_name)
    index = pc.Index(host=host)

    for i in range(0, len(chunks), batch_size):
        batch = chunks[i:i+batch_size]
        texts = [c["text"] for c in batch]
        embs = embed_texts(texts, engine)
        if not embs:
            logging.warning("Empty embedding batch; skipping.")
            continue
        vectors = []
        for c, e in zip(batch, embs):
            meta = c.get("metadata", {}).copy()
            # Keep metadata compact & serializable
            for k, v in list(meta.items()):
                if v is None:
                    meta.pop(k, None)
                elif not isinstance(v, (str, int, float, bool)):
                    meta[k] = str(v)
            vectors.append((c["id"], e, meta))
        index.upsert(vectors)

# ===============================================================
# Reporting
# ===============================================================

def build_report(segments: List[Segment], chunks: List[Dict], index_name: str) -> Dict:
    per_source: Dict[str, Dict] = {}
    for s in segments:
        key = s.locator
        rec = per_source.setdefault(key, {
            "source_type": s.source_type,
            "type": s.meta.get("type"),
            "words": 0,
            "pages": set(),  # for PDFs
            "file_path": s.meta.get("file_path"),
            "source_url": s.meta.get("source_url")
        })
        w = count_words(s.text)
        rec["words"] += w
        if "page" in s.meta:
            rec["pages"].add(s.meta["page"])

    for k, rec in per_source.items():
        if isinstance(rec.get("pages"), set):
            rec["pages"] = sorted(list(rec["pages"]))

    # chunks per locator
    chunks_per: Dict[str, int] = {}
    for c in chunks:
        loc = c.get("metadata", {}).get("locator", "unknown")
        chunks_per[loc] = chunks_per.get(loc, 0) + 1

    # merge chunks count back
    for loc, cnt in chunks_per.items():
        per_source.setdefault(loc, {}).setdefault("chunks", 0)
        per_source[loc]["chunks"] = cnt

    totals = {
        "sources": len(per_source),
        "total_words": sum(rec.get("words", 0) for rec in per_source.values()),
        "total_chunks": len(chunks)
    }

    return {
        "index_name": index_name,
        "generated_at": now_iso(),
        "totals": totals,
        "sources": per_source
    }

# ===============================================================
# Pipeline Orchestration
# ===============================================================

def run_pipeline(
    index_name: str,
    *,
    # Confluence auth + inputs
    confluence_username: str = "",
    confluence_api_token: str = "",
    confluence_pages: Optional[List[str]] = None,
    confluence_roots: Optional[List[str]] = None,
    max_pages: int = 150,
    max_depth: int = 3,
    # Coveo (optional)
    coveo_org_id: str = "",
    coveo_platform_token: str = "",
    coveo_user_email: str = "",
    coveo_labels: Optional[List[str]] = None,
    # Local files
    files: Optional[List[str]] = None,
    # Output
    report_path: str = "report.json",
):
    # 1) Secrets / SDK setup
    logging.info("Initializing secrets and SDKs...")
    pc = pinecone_config()
    embedding_engine = openai_api_config()

    # 2) Prepare URLs from Confluence inputs & Coveo
    urls: List[str] = []
    confluence_pages = confluence_pages or []
    confluence_roots = confluence_roots or []
    coveo_labels = coveo_labels or []
    files = files or []

    # Add explicit single pages
    urls.extend([u for u in confluence_pages if looks_like_url(u)])

    # Expand roots
    if confluence_roots and confluence_username and confluence_api_token:
        logging.info(f"Crawling Confluence roots (count={len(confluence_roots)})...")
        for root in confluence_roots:
            try:
                expanded = confluence_collect_descendants(root, confluence_username, confluence_api_token, max_pages=max_pages, max_depth=max_depth)
                urls.extend(expanded)
            except Exception as e:
                logging.warning(f"Root crawl failed for {root}: {e}")

    # Coveo label search → Confluence URLs
    if coveo_org_id and coveo_platform_token and coveo_labels:
        logging.info(f"Querying Coveo labels (count={len(coveo_labels)})...")
        try:
            cv = CoveoSearch(coveo_org_id, coveo_platform_token)
            user_token = cv.get_token(coveo_user_email or "user@example.com")
            for lab in coveo_labels:
                try:
                    found = cv.search_links(lab, user_token)
                    urls.extend(found)
                except Exception as e:
                    logging.warning(f"Coveo search failed for label '{lab}': {e}")
        except Exception as e:
            logging.warning(f"Coveo bootstrap failed: {e}")

    # Deduplicate URLs
    seen = set()
    deduped_urls = []
    for u in urls:
        if u not in seen:
            seen.add(u)
            deduped_urls.append(u)

    # 3) Extract segments
    segments: List[Segment] = []

    # Confluence extraction
    if deduped_urls and confluence_username and confluence_api_token:
        logging.info(f"Extracting Confluence pages: {len(deduped_urls)}")
        segments.extend(extract_confluence_pages(deduped_urls, confluence_username, confluence_api_token))
    elif deduped_urls and not (confluence_username and confluence_api_token):
        logging.warning("Confluence URLs provided but no credentials; skipping Confluence extraction.")

    # Files extraction
    for f in files:
        if not os.path.isfile(f):
            logging.warning(f"File not found: {f}")
            continue
        segments.extend(extract_file_segments(f))

    if not segments:
        logging.warning("No segments extracted — nothing to index.")
        report = build_report([], [], index_name)
        with open(report_path, "w", encoding="utf-8") as f:
            json.dump(report, f, indent=2)
        logging.info(f"Report written: {report_path}")
        return

    # 4) Chunking
    logging.info(f"Chunking {len(segments)} segments...")
    chunks = chunk_segments(segments)
    logging.info(f"Total chunks: {len(chunks)}")

    # 5) Embeddings + Pinecone Upsert (index reuse-or-create)
    logging.info(f"Upserting into Pinecone index '{index_name}'...")
    upsert_chunks(pc, index_name, chunks, embedding_engine)

    # 6) Report
    report = build_report(segments, chunks, index_name)
    with open(report_path, "w", encoding="utf-8") as f:
        json.dump(report, f, indent=2)
    logging.info(f"Report written: {report_path}")
    logging.info("Pipeline completed successfully.")


# ===============================================================
# Main: read non-secret config.json and run
# ===============================================================

if __name__ == "__main__":
    cfg_path = os.path.join(os.path.dirname(__file__), "config.json")
    if not os.path.isfile(cfg_path):
        example = {
            "index_name": "test-team-index",

            "confluence_username": "you@company.com",
            "confluence_api_token": "ATLTOKEN_xxx",
            "confluence_pages": [
                "https://your.atlassian.net/wiki/pages/123456789"
            ],
            "confluence_roots": [
                "https://your.atlassian.net/wiki/pages/987654321"
            ],
            "max_pages": 120,
            "max_depth": 3,

            "coveo_org_id": "",
            "coveo_platform_token": "",
            "coveo_user_email": "",
            "coveo_labels": [],

            "files": [
                # "docs/policy.pdf",
                # "docs/runbook.docx",
                # "docs/notes.txt",
                # "docs/data.xlsx",
                # "docs/slides.pptx"
            ],

            "report_path": "report.json"
        }
        print("No config.json found. Here is an example you can save next to pipeline.py:")
        print(json.dumps(example, indent=2))
        raise SystemExit(1)

    with open(cfg_path, "r", encoding="utf-8") as f:
        cfg = json.load(f)

    # sensible defaults
    cfg.setdefault("confluence_pages", [])
    cfg.setdefault("confluence_roots", [])
    cfg.setdefault("coveo_labels", [])
    cfg.setdefault("files", [])
    cfg.setdefault("max_pages", 150)
    cfg.setdefault("max_depth", 3)
    cfg.setdefault("report_path", "report.json")

    run_pipeline(**cfg)
